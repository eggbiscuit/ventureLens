"""
BP解析Agent
解析商业计划书提取结构化信息
"""

import asyncio
from typing import Dict, Any, List, Optional
import logging
from datetime import datetime
import os

from agents.base import BaseAgent
from state import VentureLensState

logger = logging.getLogger(__name__)


class BPParserAgent(BaseAgent):
    """BP解析Agent"""
    
    def __init__(self, config: Dict[str, Any]):
        super().__init__("bp_parser", config)
        
    async def _execute(self, state: VentureLensState) -> VentureLensState:
        """执行BP解析"""
        
        bp_file_path = state.get("bp_file_path")
        if not bp_file_path or not os.path.exists(bp_file_path):
            logger.info("No BP file provided or file does not exist, skipping BP parsing")
            return state
        
        company_name = state["company_name"]
        
        try:
            # 解析BP文件
            bp_content = await self._extract_bp_content(bp_file_path)
            
            if not bp_content:
                logger.warning("Failed to extract content from BP file")
                return state
            
            # 使用LLM结构化BP内容
            structured_data = await self._structure_bp_data(company_name, bp_content)
            
            # 更新状态
            state["bp_extracted_data"] = structured_data
            
            logger.info(f"Successfully parsed BP for {company_name}")
            
        except Exception as e:
            logger.error(f"Error parsing BP: {e}")
            state["bp_extracted_data"] = {"error": str(e)}
        
        return state
    
    async def _extract_bp_content(self, file_path: str) -> str:
        """从BP文件中提取文本内容"""
        
        file_extension = os.path.splitext(file_path)[1].lower()
        
        try:
            if file_extension == '.pdf':
                return await self._extract_from_pdf(file_path)
            elif file_extension in ['.ppt', '.pptx']:
                return await self._extract_from_pptx(file_path)
            elif file_extension in ['.doc', '.docx']:
                return await self._extract_from_docx(file_path)
            else:
                logger.warning(f"Unsupported file format: {file_extension}")
                return ""
                
        except Exception as e:
            logger.error(f"Error extracting content from {file_path}: {e}")
            return ""
    
    async def _extract_from_pdf(self, file_path: str) -> str:
        """从PDF中提取文本"""
        try:
            import pypdf
            
            text_content = []
            with open(file_path, 'rb') as file:
                pdf_reader = pypdf.PdfReader(file)
                for page in pdf_reader.pages:
                    text_content.append(page.extract_text())
            
            return "\n".join(text_content)
            
        except ImportError:
            logger.error("pypdf not installed, cannot parse PDF files")
            return ""
        except Exception as e:
            logger.error(f"Error parsing PDF: {e}")
            return ""
    
    async def _extract_from_pptx(self, file_path: str) -> str:
        """从PPTX中提取文本"""
        try:
            from pptx import Presentation
            
            text_content = []
            prs = Presentation(file_path)
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content.append(shape.text)
            
            return "\n".join(text_content)
            
        except ImportError:
            logger.error("python-pptx not installed, cannot parse PPTX files")
            return ""
        except Exception as e:
            logger.error(f"Error parsing PPTX: {e}")
            return ""
    
    async def _extract_from_docx(self, file_path: str) -> str:
        """从DOCX中提取文本"""
        try:
            from docx import Document
            
            doc = Document(file_path)
            text_content = [paragraph.text for paragraph in doc.paragraphs]
            
            return "\n".join(text_content)
            
        except ImportError:
            logger.error("python-docx not installed, cannot parse DOCX files")
            return ""
        except Exception as e:
            logger.error(f"Error parsing DOCX: {e}")
            return ""
    
    async def _structure_bp_data(self, company_name: str, bp_content: str) -> Dict[str, Any]:
        """使用LLM结构化BP内容"""
        
        prompt = f"""
请分析以下商业计划书内容，提取关键信息并结构化输出：

公司名称: {company_name}

BP内容:
{bp_content[:8000]}  # 限制长度避免token过多

请以以下JSON格式输出关键信息：
{{
    "company_info": {{
        "name": "公司名称",
        "founding_date": "成立时间",
        "business_model": "商业模式描述",
        "core_product": "核心产品/服务",
        "target_market": "目标市场"
    }},
    "team_info": {{
        "founders": ["创始人1", "创始人2"],
        "key_members": ["关键成员1", "关键成员2"],
        "team_size": "团队规模"
    }},
    "financial_info": {{
        "funding_sought": "寻求融资金额",
        "valuation": "估值",
        "revenue": "营收情况",
        "projections": "财务预测",
        "use_of_funds": "资金用途"
    }},
    "market_info": {{
        "industry": "所属行业",
        "market_size": "市场规模",
        "competitors": ["竞争对手1", "竞争对手2"],
        "competitive_advantage": "竞争优势"
    }},
    "business_metrics": {{
        "key_metrics": ["关键指标1", "关键指标2"],
        "growth_metrics": "增长指标",
        "customer_metrics": "客户指标"
    }},
    "risks_and_challenges": ["风险点1", "风险点2"],
    "milestones": ["里程碑1", "里程碑2"]
}}
"""
        
        try:
            response = await self.llm_service._call_llm(prompt)
            return self.llm_service._parse_analysis_response(response)
        except Exception as e:
            logger.error(f"Error structuring BP data: {e}")
            return {"error": str(e), "raw_content": bp_content[:1000]}
